
# Handles text extraction from PDF, Word, and PowerPoint files.
# Unified extract_text function determines the file type and calls the appropriate extraction method.


import PyPDF2
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(pdf_path):
    text = ""
    with open(pdf_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(docx_path):
    text = ""
    doc = Document(docx_path)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(pptx_path):
    text = ""
    presentation = Presentation(pptx_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
    return text

def extract_text(file_path, file_type):
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type == "docx":
        return extract_text_from_docx(file_path)
    elif file_type == "pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file type")